import os.path

from django.shortcuts import render
from django import forms
from django.http import HttpResponse
from pptx import Presentation
from django.core.files import File
from ppt.models import PptModel
from spire.presentation import Presentation as pv
import pyttsx3
from moviepy.editor import concatenate, ImageClip
import moviepy.editor as mpe
import os
from pptx import Presentation
from gtts import gTTS
from moviepy.editor import concatenate_videoclips, ImageClip, AudioFileClip
from PIL import Image
from io import BytesIO
import comtypes.client


class PptModelForm(forms.ModelForm):
    class Meta:
        model = PptModel
        fields = ['pptFile']


def index2(request):
    if request.method == 'POST':
        form = PptModelForm(request.POST, request.FILES)
        if form.is_valid():
            obj = form.save()
            c = ''
            text_speech = pyttsx3.init()
            prs = Presentation(obj.pptFile)
            for slide in prs.slides:
                print('12')
                if slide.notes_slide:
                    comments = slide.notes_slide.notes_text_frame.text
                    print(comments)
                    c += comments
                    print(dir(slide))
                    if comments:
                        rate = 150
                        text_speech.setProperty('voice', 'HKEY_LOCAL_MACHINE\SOFTWARE\Microsoft\Speech\Voices\Tokens\TTS_MS_EN-\
                                           US_DAVID_11.0')
                        text_speech.setProperty('rate', rate)
            print(c)
            text_speech.save_to_file(c, 'media/abc.Mp3')
            text_speech.runAndWait()
            presentation = pv()
            presentation.LoadFromFile("media\\" + str(obj.pptFile))
            fileNames = []
            for i, slide in enumerate(presentation.Slides):
                fileNames.append("ToImage_" + str(i) + ".png")
                img = slide.SaveAsImage()
                img.Save(fileNames[-1])
                img.Dispose()
            presentation.Dispose()
            print([os.path.join(os.path.relpath('.'), i) for i in fileNames])
            video = concatenate([ImageClip(os.path.join(os.path.relpath('.'), i)).set_duration(1) for i in fileNames],
                                method="compose")
            audio_background = mpe.AudioFileClip('media/abc.Mp3')
            final_clip = video.set_audio(audio_background)
            final_clip.write_videofile("media/output.mp4", fps=24)
            return render(request, 'index.html', context={'audio': 'output.mp4'})
    return render(request, 'index.html')

def save_slide_as_image(ppt_file_path, slide_index, output_image_path):
    # Use comtypes to export slide as image
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1
    presentation = powerpoint.Presentations.Open(os.path.abspath(ppt_file_path))
    slide = presentation.Slides[slide_index + 1]  # Slide index is 1-based in PowerPoint
    slide.Export(os.path.abspath(output_image_path), "PNG")
    presentation.Close()
    powerpoint.Quit()
def index(request):
    output_audio_path='output_audio.mp3'
    output_video_path='output_video.mp4'
    if request.method == 'POST':
        form = PptModelForm(request.POST, request.FILES)
        if form.is_valid():
            obj = form.save()
            c = ''
            # text_speech = pyttsx3.init()
            print(obj.pptFile)
            image_files = []
            presentation = pv()
            presentation.LoadFromFile("media\\" + str(obj.pptFile))
            for i, slide in enumerate(presentation.Slides):
                print(slide)
                image_files.append("ToImage_" + str(i) + ".png")
                img = slide.SaveAsImage()
                img.Save(image_files[-1])
                img.Dispose()
            presentation.Dispose()
            prs = Presentation(obj.pptFile)
            media_dir= os.path.abspath('media')
            ppt_file_path=obj.pptFile

            audio_files = []

            for i, slide in enumerate(prs.slides):
                # Save slide as image
                slide_image_path = os.path.join(media_dir, f"slide_{i}.png")
                # save_slide_as_image(ppt_file_path, i, slide_image_path)
                # image_files.append(slide_image_path)

                # Extract notes text
                notes_text = ''
                if slide.notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()

                # Convert notes text to speech and save as audio file
                if notes_text:
                    audio_file_path = os.path.join(media_dir, f"audio_{i}.mp3")
                    tts = gTTS(text=notes_text, lang='en')
                    tts.save(audio_file_path)
                    audio_files.append(audio_file_path)
                else:
                    audio_files.append(None)  # Append None if there's no audio

            # Create video clips for each slide
            video_clips = []
            for img_file, audio_file in zip(image_files, audio_files):
                img_clip = ImageClip(img_file)
                if audio_file:
                    audio_clip = AudioFileClip(audio_file)
                    img_clip = img_clip.set_duration(audio_clip.duration)
                    video_clip = img_clip.set_audio(audio_clip)
                else:
                    img_clip = img_clip.set_duration(2)  # Default duration for silent slides
                    video_clip = img_clip
                video_clips.append(video_clip)

            # Concatenate video clips
            final_video = concatenate_videoclips(video_clips, method="compose")
            final_video.write_videofile(output_video_path, fps=24)

            # Clean up intermediate files
            for img in image_files:
                os.remove(img)
            for audio in audio_files:
                if audio:
                    os.remove(audio)
            print(output_video_path)
            return render(request, 'index.html', context={'audio': f'{output_video_path}.mp4'})
    return render(request, 'index.html')
